import io
import os
import zipfile
from xml.etree import ElementTree as ET

import pandas as pd
import streamlit as st

from crawl import client
from utils.helpers import fill_missing_cells, remove_code_fence


def extract_text_from_pdf(file):
    import io
    import pdfplumber
    from pdfminer.layout import LAParams
    import pandas as pd

    file_bytes = file.read()
    pdf_io = io.BytesIO(file_bytes)
    text = ""
    laparams = LAParams(char_margin=2.0, line_margin=0.5, word_margin=0.1, boxes_flow=0.5)
    with pdfplumber.open(pdf_io, laparams=laparams.__dict__) as pdf:
        for page in pdf.pages:
            # 추출된 일반 텍스트
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"

            # 페이지 내 테이블 추출
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    # 병합된 셀(누락된 값) 보완: 이전 행의 값으로 채워 넣기
                    table = fill_missing_cells(table)
                    # 테이블을 Markdown으로 변환 (첫 행을 헤더로 사용)
                    try:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        markdown_table = df.to_markdown(index=False)
                    except Exception:
                        # DataFrame 변환에 실패하면 단순 문자열 결합 방식 사용
                        markdown_table = "\n".join([" | ".join(map(str, row)) for row in table])
                    text += "\n" + markdown_table + "\n\n"
    return text


def extract_text_from_hwpx(file):
    try:
        file_bytes = file.read()
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            names = z.namelist()
            candidate = next((name for name in names if "section0.xml" in name.lower()), None)
            if not candidate:
                st.error("hwpx 파일에서 section XML을 찾을 수 없습니다.")
                return ""
            xml_content = z.read(candidate)
    except Exception as e:
        st.error(f"hwpx 파일 처리 오류: {e}")
        return ""
    try:
        tree = ET.fromstring(xml_content)
    except Exception as e:
        st.error(f"XML 파싱 오류: {e}")
        return ""
    texts = [elem.text.strip() for elem in tree.iter() if elem.text]
    return " ".join(texts).strip()


def extract_text_from_txt(file):
    try:
        content = file.read()
        if isinstance(content, bytes):
            content = content.decode("utf-8")
    except Exception as e:
        st.error(f"TXT 파일 추출 오류: {e}")
        content = ""
    return content


def extract_text_from_pptx(file):
    try:
        from pptx import Presentation
    except ImportError:
        st.error("pptx 파일 처리를 위해 python-pptx 라이브러리가 필요합니다. pip install python-pptx")
        return ""
    try:
        file.seek(0)
        ppt = Presentation(io.BytesIO(file.read()))
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"PPTX 파일 추출 오류: {e}")
        return ""


def extract_text_from_xlsx(file):
    try:
        file_bytes = file.read()
        from openpyxl import load_workbook
        wb = load_workbook(filename=io.BytesIO(file_bytes), data_only=True)
        result_text = ""
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            # 병합된 셀 처리: 각 병합 영역에 대해 unmerge 후, 모든 셀에 왼쪽 상단 값 채우기
            merged_ranges = list(ws.merged_cells.ranges)
            for merged_range in merged_ranges:
                top_left_value = ws.cell(merged_range.min_row, merged_range.min_col).value
                # 병합 영역 해제
                ws.unmerge_cells(range_string=str(merged_range))
                # 해제된 각 셀에 값을 채워 넣기
                for row in range(merged_range.min_row, merged_range.max_row + 1):
                    for col in range(merged_range.min_col, merged_range.max_col + 1):
                        ws.cell(row=row, column=col).value = top_left_value
            # 시트의 데이터를 리스트 형태로 추출 (첫 행은 컬럼명)
            data = ws.values
            columns = next(data)
            df = pd.DataFrame(data, columns=columns)
            try:
                markdown_table = df.to_markdown(index=False)
            except Exception:
                markdown_table = df.to_string(index=False)
            result_text += f"Sheet: {sheet_name}\n{markdown_table}\n\n"
        return result_text
    except Exception as e:
        st.error(f"XLSX 파일 추출 오류: {e}")
        return ""


def extract_text_from_file(file):
    ext = os.path.splitext(file.name)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file)
    elif ext == ".hwp":
        st.error("HWP 파일 추출 기능이 구현되어 있지 않습니다.")
        return ""
    elif ext == ".hwpx":
        return extract_text_from_hwpx(file)
    elif ext == ".txt":
        return extract_text_from_txt(file)
    elif ext == ".pptx":
        return extract_text_from_pptx(file)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file)
    else:
        st.error("지원되지 않는 파일 형식입니다. (PDF, HWP, HWPX, TXT, PPTX, XLSX 파일만 지원)")
        return ""


def split_text_with_overlap(text, chunk_size, overlap):
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunks.append(" ".join(words[start:end]))
        start = end - overlap
    return chunks


def convert_chunk_to_markdown(current_chunk, previous_md=None):
    context = ""
    if previous_md:
        prev_words = previous_md.split()
        context = " ".join(prev_words[-100:]) + "\n\n"
    prompt = (
        "**현재청크**를 잘 구조화된 Markdown 문서로 변환해줘. "
        "Markdown 문서의 문맥을 분석하여 어색한 부분들은 모두 삭제해줘. "
        "**이전 마크다운 청크**는 포함하지 않고 **현재청크**만 정제해서 만들어줘. "
        f"**이전 마크다운 청크**: {context}\n\n"
        f"**현재청크**: {current_chunk}"
    )
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": "너는 텍스트 변환 전문가야. 주어진 청크의 모든 세부 정보를 보존하면서 Markdown 문서로 변환해."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.3
    )
    markdown_text = response.choices[0].message.content
    return remove_code_fence(markdown_text)
